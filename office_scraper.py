import asyncio
import json
import re
import tempfile
import os
import requests
from urllib.parse import urlparse, urljoin
from bs4 import BeautifulSoup
from pathlib import Path
from playwright.async_api import async_playwright
import fitz  # PyMuPDF

# Office document libraries
import docx
import openpyxl
from pptx import Presentation


async def google_search(query, num_results=10):
    """Perform a Google search and return a list of result URLs."""
    urls = []
    print(f"Searching Google for: {query}")
    
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        context = await browser.new_context(
            user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        )
        page = await context.new_page()

        # Navigate to Google and perform search
        await page.goto(f"https://www.google.com/search?q={query}", timeout=60000)
        
        # Wait for search results to load
        await page.wait_for_selector("div#search", timeout=10000)
        
        # Extract result links
        result_links = await page.query_selector_all("div.g a[href^='http']")
        
        for link in result_links:
            href = await link.get_attribute("href")
            if href and href.startswith("http"):
                # Filter out Google's own domains
                parsed_url = urlparse(href)
                if "google" not in parsed_url.netloc:
                    urls.append(href)
                    print(f"Found URL: {href}")
            
            if len(urls) >= num_results:
                break
                
        # If we didn't get enough results with the first selector, try another common one
        if len(urls) < num_results:
            elements = await page.query_selector_all("a")
            for e in elements:
                href = await e.get_attribute("href")
                if href and href.startswith("/url?q="):
                    clean_url = href.split("/url?q=")[1].split("&sa=")[0]
                    if not clean_url.startswith("http"):
                        continue
                    
                    parsed_url = urlparse(clean_url)
                    if "google" not in parsed_url.netloc and clean_url not in urls:
                        urls.append(clean_url)
                        print(f"Found URL: {clean_url}")
                
                if len(urls) >= num_results:
                    break

        await browser.close()
    
    return list(set(urls))


def is_valid_url(url):
    """Check if a URL is valid."""
    try:
        result = urlparse(url)
        return all([result.scheme, result.netloc])
    except:
        return False


def normalize_url(base_url, url):
    """Convert relative URLs to absolute URLs."""
    if is_valid_url(url):
        return url
    return urljoin(base_url, url)


def crawl_page(url):
    """Extract text and downloadable links from a webpage."""
    print(f"Crawling {url}")
    text = ""
    downloadables = []
    
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        response = requests.get(url, timeout=15, headers=headers)
        response.raise_for_status()
        
        content_type = response.headers.get('Content-Type', '').lower()
        
        # If the URL is a direct file
        if any(ct in content_type for ct in ['application/pdf', 'application/vnd.openxmlformats-officedocument', 'application/msword']):
            return "", [url]
            
        # If the URL is a direct text file
        if 'text/plain' in content_type:
            return response.text, []
        
        # Parse HTML content
        soup = BeautifulSoup(response.text, "html.parser")
        
        # Extract text content
        for script in soup(["script", "style"]):
            script.extract()
        text = soup.get_text(separator=" ", strip=True)
        
        # Extract links to downloadable files
        links = []
        for a in soup.find_all('a', href=True):
            href = a['href']
            full_url = normalize_url(url, href)
            if full_url and is_valid_url(full_url):
                links.append(full_url)
        
        # Filter for downloadable files
        downloadables = [
            link for link in links 
            if any(link.lower().endswith(ext) for ext in ['.pdf', '.txt', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'])
        ]
        
        return text, downloadables
        
    except Exception as e:
        print(f"Error crawling {url}: {e}")
        return "", []


def extract_pdf_text(url):
    """Download and extract text from a PDF file."""
    print(f"Extracting PDF from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            doc = fitz.open(temp_path)
            text = ""
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text += page.get_text()
            doc.close()
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing PDF file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading PDF from {url}: {e}")
        return ""


def extract_docx_text(url):
    """Download and extract text from a DOCX file."""
    print(f"Extracting DOCX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            doc = docx.Document(temp_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing DOCX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading DOCX from {url}: {e}")
        return ""


def extract_xlsx_text(url):
    """Download and extract text from an XLSX file."""
    print(f"Extracting XLSX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            workbook = openpyxl.load_workbook(temp_path, data_only=True)
            text = ""
            
            # Extract text from each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows():
                    row_text = " | ".join(str(cell.value) if cell.value is not None else "" for cell in row)
                    if row_text.strip():
                        text += row_text + "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing XLSX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading XLSX from {url}: {e}")
        return ""


def extract_pptx_text(url):
    """Download and extract text from a PPTX file."""
    print(f"Extracting PPTX from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, stream=True, timeout=30, headers=headers)
        r.raise_for_status()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as f:
            f.write(r.content)
            temp_path = f.name
            
        try:
            presentation = Presentation(temp_path)
            text = ""
            
            for i, slide in enumerate(presentation.slides):
                text += f"Slide {i+1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
            
            os.unlink(temp_path)  # Remove the temporary file
            return text
        except Exception as e:
            print(f"Error processing PPTX file {url}: {e}")
            os.unlink(temp_path)  # Ensure the temp file is removed
            return ""
            
    except Exception as e:
        print(f"Error downloading PPTX from {url}: {e}")
        return ""


def extract_txt_text(url):
    """Download and extract text from a text file."""
    print(f"Extracting TXT from {url}")
    try:
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/96.0.4664.110 Safari/537.36"
        }
        r = requests.get(url, timeout=15, headers=headers)
        r.raise_for_status()
        return r.text
    except Exception as e:
        print(f"Error extracting TXT from {url}: {e}")
        return ""


def clean_text(text):
    """Clean and normalize extracted text."""
    # Remove URLs
    text = re.sub(r'https?://\S+', '', text)
    
    # Remove email addresses
    text = re.sub(r'\S+@\S+\.\S+', '', text)
    
    # Convert None values to empty strings
    text = re.sub(r'None', '', text)
    
    # Remove extra whitespace, newlines, tabs
    text = re.sub(r'\s+', ' ', text)
    
    # Remove very long words (likely garbage)
    text = re.sub(r'\S{50,}', '', text)
    
    return text.strip()


def save_json(data, filename="output.json"):
    """Save the results to a JSON file."""
    with open(filename, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"Data saved to {filename}")


async def extract_document_text(url):
    """Extract text from document URLs based on file extension."""
    url_lower = url.lower()
    
    if url_lower.endswith('.pdf'):
        return extract_pdf_text(url)
    elif url_lower.endswith('.txt'):
        return extract_txt_text(url)
    elif url_lower.endswith('.docx'):
        return extract_docx_text(url)
    elif url_lower.endswith('.xlsx'):
        return extract_xlsx_text(url)
    elif url_lower.endswith('.pptx'):
        return extract_pptx_text(url)
    elif url_lower.endswith('.doc'):
        print(f"Legacy DOC format not supported: {url}")
        return ""
    elif url_lower.endswith('.xls'):
        print(f"Legacy XLS format not supported: {url}")
        return ""
    elif url_lower.endswith('.ppt'):
        print(f"Legacy PPT format not supported: {url}")
        return ""
    else:
        return ""


async def process_url(url):
    """Process a single URL and extract all relevant information."""
    page_text, downloadables = crawl_page(url)
    
    # Process downloadable documents
    document_texts = []
    for doc_url in downloadables:
        doc_text = await extract_document_text(doc_url)
        if doc_text:
            document_texts.append(doc_text)
    
    # Combine all texts
    combined_text = page_text
    if document_texts:
        combined_text += "\n" + "\n".join(document_texts)
    
    # Clean the combined text
    cleaned_text = clean_text(combined_text)
    
    return {
        "url": url,
        "content": cleaned_text[:10000]  # Limit to 10k chars to prevent huge outputs
    }


async def main(query, num_results=10, output_file="output.json"):
    """Main function to run the scraper."""
    print(f"Starting scraper for query: {query}")
    
    # Get URLs from Google search
    urls = await google_search(query, num_results)
    print(f"Found {len(urls)} URLs")
    
    # Process each URL
    tasks = [process_url(url) for url in urls]
    results = await asyncio.gather(*tasks)
    
    # Remove empty results
    results = [r for r in results if r["content"]]
    
    # Save results
    save_json(results, output_file)
    print(f"Scraping completed. Found content from {len(results)} pages.")
    
    return results


if __name__ == "__main__":
    import sys
    import argparse
    
    parser = argparse.ArgumentParser(description="Office Document Scraper")
    parser.add_argument("query", nargs="*", help="Search query")
    parser.add_argument("--results", type=int, default=10, help="Number of search results to process")
    parser.add_argument("--output", default="output.json", help="Output JSON file")
    
    args = parser.parse_args()
    
    query = " ".join(args.query) if args.query else "site:gov climate change report"
    
    print(f"Starting scraper with query: {query}")
    asyncio.run(main(query, args.results, args.output)) 